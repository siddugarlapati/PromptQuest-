import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

def create_pitch_deck():
    prs = Presentation()
    
    # 16:9 Aspect Ratio
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Define colors
    bg_color = RGBColor(15, 23, 42) # #0f172a (Tailwind slate-900)
    text_color = RGBColor(255, 255, 255)
    accent_color = RGBColor(236, 72, 153) # #ec4899 (Pink-500)
    
    logo_path = "/Users/garlapati/Coding/AI-simulation/frontend/public/au-logo.png"

    def add_slide(title_text, subtitle_text="", content_bullets=None, include_image_placeholder=False, placeholder_text="[Insert GUI Screenshot Here]"):
        # Use blank layout
        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Set dark background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color
        
        # Add Logo (top right)
        if os.path.exists(logo_path):
            slide.shapes.add_picture(logo_path, Inches(11.5), Inches(0.5), height=Inches(0.8))
        
        # Add Title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(44)
        p.font.color.rgb = accent_color
        
        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.size = Pt(24)
            p2.font.color.rgb = text_color
            
        # Add Bullets
        if content_bullets:
            content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6 if include_image_placeholder else 11), Inches(4.5))
            tf_content = content_box.text_frame
            tf_content.word_wrap = True
            
            for i, bullet in enumerate(content_bullets):
                p = tf_content.add_paragraph() if i > 0 else tf_content.paragraphs[0]
                p.text = "• " + bullet
                p.font.size = Pt(22)
                p.font.color.rgb = text_color
                p.space_after = Pt(14)
                
        # Add Image Placeholder
        if include_image_placeholder:
            shape = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, 
                Inches(7.5), Inches(2), Inches(5), Inches(4.5)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(30, 41, 59)
            shape.line.color.rgb = accent_color
            shape.line.width = Pt(2)
            
            # Placeholder text
            tf_ph = shape.text_frame
            tf_ph.text = placeholder_text
            tf_ph.paragraphs[0].alignment = PP_ALIGN.CENTER
            tf_ph.paragraphs[0].font.size = Pt(20)
            tf_ph.paragraphs[0].font.color.rgb = RGBColor(148, 163, 184)
            
        return slide

    # SLIDE 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    slide1.background.fill.solid()
    slide1.background.fill.fore_color.rgb = bg_color
    if os.path.exists(logo_path):
        slide1.shapes.add_picture(logo_path, Inches(5.666), Inches(1.5), height=Inches(1.5))
    
    tb = slide1.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "PromptQuest"
    p.font.bold = True
    p.font.size = Pt(64)
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "The Next Generation of AI Education at Anurag University"
    p2.font.size = Pt(28)
    p2.font.color.rgb = text_color
    p2.alignment = PP_ALIGN.CENTER

    # SLIDE 2: The Problem
    add_slide(
        "The AI Education Gap",
        "Moving from Consumers to Engineers",
        [
            "The Black Box Problem: Students view AI as magic. They lack understanding of Transformers, Tokens, and Embeddings.",
            "The Cost Barrier: Real API keys (OpenAI, Anthropic) are too expensive to scale across thousands of university students.",
            "The Engagement Drop: Academic AI theory is dense; it demands interactive visualization."
        ]
    )

    # SLIDE 3: The Solution
    add_slide(
        "Introducing PromptQuest",
        "A 100% Free, Local, Gamified AI Sandbox",
        [
            "Zero Cost: Runs entirely on local, open-source models (Llama 3.2 via Ollama).",
            "Maximum Privacy: Student data and prompts never leave the local machine.",
            "Gamified Learning: Complex math is turned into playable mini-games providing visual and immediate feedback.",
            "Hands-on Engineering: Students actually build apps."
        ],
        include_image_placeholder=True,
        placeholder_text="[Insert PromptQuest Dashboard UI Screenshot]"
    )

    # SLIDE 4: curriculum
    add_slide(
        "Interactive Curriculum",
        "14 Worlds Covering Core AI Mechanics",
        [
            "Worlds 1-3 (Basics): Pattern recognition, text prediction, and Tokenization.",
            "Worlds 4-7 (The Brain): Visualizing Transformers, Vector Embeddings, and the Context Window.",
            "Worlds 8-10 (Application): Prompt Engineering frameworks and combatting AI Hallucinations."
        ],
        include_image_placeholder=True,
        placeholder_text="[Insert 14-World Hub UI Screenshot]"
    )

    # SLIDE 5: Arcade
    add_slide(
        "The LLM Arcade (Patentable Innovation)",
        "Gamifying Neural Engine Mechanics",
        [
            "The Autoregressor: Race a timer to predict the highest-probability Next Token.",
            "The Context Squeezer: Compress bloated prompts down to strict token limits.",
            "The Attention Mapper: Map Attention Weights connecting pronouns to nouns.",
            "Infinite Replayability: Games are NOT hardcoded; they are generated dynamically on-the-fly by the local LLM."
        ],
        include_image_placeholder=True,
        placeholder_text="[Insert LLM Arcade UI Scene Screenshot]"
    )
    
    # SLIDE 6: UI
    add_slide(
        "Vibe Coding & Generative UI",
        "The Future of Software Development",
        [
            "Students learn to build apps like Cursor and v0.",
            "They type a prompt, watch the local AI stream React code, and see the application render live in-browser via Sandpack.",
            "Real-time transition from writing code to directing AI logic."
        ],
        include_image_placeholder=True,
        placeholder_text="[Insert Vibe Coding Sandpack UI Screenshot]"
    )

    # SLIDE 7: RAG
    add_slide(
        "Enterprise RAG Pipelines",
        "Teaching Industry-Standard Architecture",
        [
            "Students implement Retrieval-Augmented Generation (RAG).",
            "They upload PDFs, watch the Vector DB (ChromaDB) chunk and embed text.",
            "Force the AI to answer strictly based on secured, local data.",
            "The #1 most requested skill by tech employers today."
        ],
        include_image_placeholder=True,
        placeholder_text="[Insert RAG Pipeline UI Screenshot]"
    )
    
    # SLIDE 8: Roadmap
    add_slide(
        "Implementation Roadmap",
        "How We Scale This Across Campus",
        [
            "Phase 1: Lab Integration. Install Ollama & PromptQuest on Campus Computer Labs.",
            "Phase 2: Curriculum Pilot. Use the 14 Worlds as lab assignments for Intro CS classes.",
            "Phase 3: Hackathons. Host university-wide competitions using the LLM Arcade."
        ]
    )

    # SLIDE 9: Conclusion
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    slide9.background.fill.solid()
    slide9.background.fill.fore_color.rgb = bg_color
    if os.path.exists(logo_path):
        slide9.shapes.add_picture(logo_path, Inches(5.666), Inches(1.5), height=Inches(1.5))
        
    tb = slide9.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11.33), Inches(2))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = "Pioneering AI Literacy"
    p.font.bold = True
    p.font.size = Pt(54)
    p.font.color.rgb = accent_color
    p.alignment = PP_ALIGN.CENTER
    
    p2 = tf.add_paragraph()
    p2.text = "Not just surviving the AI revolution — leading it."
    p2.font.size = Pt(24)
    p2.font.color.rgb = text_color
    p2.alignment = PP_ALIGN.CENTER

    output_path = "/Users/garlapati/Desktop/PromptQuest_VC_Pitch.pptx"
    try:
        prs.save(output_path)
        print(f"Presentation generated at: {output_path}")
    except Exception as e:
        print(f"Failed to save to Desktop. Saving to local dir.")
        prs.save("PromptQuest_VC_Pitch.pptx")
        print("Presentation generated locally.")

if __name__ == '__main__':
    create_pitch_deck()
